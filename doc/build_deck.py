"""Build doc/dutopia-handbook.pptx from handbook.md content.

Run: python doc/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Palette
NAVY = RGBColor(0x0F, 0x1E, 0x3A)
TEAL = RGBColor(0x13, 0x8D, 0x90)
AMBER = RGBColor(0xE3, 0x8A, 0x1C)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
GREY = RGBColor(0x5A, 0x63, 0x73)
DARK = RGBColor(0x1A, 0x1F, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = head
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = NAVY
        if body:
            r2 = p.add_run()
            r2.text = "  " + body
            r2.font.name = font
            r2.font.size = Pt(size - 1)
            r2.font.color.rgb = color
    return tb


def add_code(slide, x, y, w, h, code, size=12):
    shp = add_rect(slide, x, y, w, h, DARK)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                   w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xD4, 0xE4, 0xFF)
    return shp


def header(slide, title, eyebrow=None):
    # Top bar
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.45), Inches(0.12), Inches(10), Inches(0.35),
             "DUTOPIA  .  TECHNICAL HANDBOOK", size=11, bold=True,
             color=WHITE)
    add_text(slide, Inches(10.8), Inches(0.12), Inches(2.2), Inches(0.35),
             f"{slide_num[0]} / {total_slides}", size=11, color=WHITE,
             align=PP_ALIGN.RIGHT)
    # Title block
    if eyebrow:
        add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.3),
                 eyebrow.upper(), size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    # Divider
    add_rect(slide, Inches(0.5), Inches(1.75), Inches(0.8), Inches(0.05),
             AMBER)


def footer(slide):
    add_text(slide, Inches(0.45), Inches(7.1), Inches(10), Inches(0.3),
             "Rust + SvelteKit  .  github.com/sganis/dutopia",
             size=9, color=GREY)


slide_num = [0]
total_slides = 14  # updated later


def new_slide():
    slide_num[0] += 1
    s = prs.slides.add_slide(BLANK)
    return s


# ======================================================================
# 1. Cover
# ======================================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
# Accent block
add_rect(s, 0, Inches(3.2), Inches(0.3), Inches(1.2), AMBER)
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.5),
         "TECHNICAL HANDBOOK  .  v4", size=14, bold=True, color=TEAL)
add_text(s, Inches(0.7), Inches(2.8), Inches(12), Inches(1.2),
         "Dutopia", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.7),
         "High-scale filesystem analytics in Rust",
         size=24, color=LIGHT)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.5),
         "Tested to 1B+ files  .  30 PB storage  .  UTF-8 clean end to end",
         size=14, color=RGBColor(0xB8, 0xC4, 0xDC))
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "For the technical team", size=12, color=TEAL, bold=True)

# ======================================================================
# 2. Pipeline overview
# ======================================================================
s = new_slide()
header(s, "The pipeline", "architecture at a glance")

# Five stages as connected boxes
stages = [
    ("duscan", "scan FS", "raw CSV / .zst"),
    ("dusum", "rollup", "sum.csv"),
    ("dudb", "load", "SQLite"),
    ("duapi", "serve", "REST + SPA"),
]
x0 = Inches(0.6)
y0 = Inches(2.3)
bw = Inches(2.6)
bh = Inches(1.6)
gap = Inches(0.4)

for i, (name, verb, out) in enumerate(stages):
    bx = x0 + (bw + gap) * i
    add_rect(s, bx, y0, bw, bh, WHITE, line=NAVY)
    add_rect(s, bx, y0, bw, Inches(0.45), NAVY)
    add_text(s, bx, y0 + Inches(0.05), bw, Inches(0.4),
             name, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, bx, y0 + Inches(0.55), bw, Inches(0.4),
             verb, size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, bx, y0 + Inches(0.95), bw, Inches(0.5),
             out, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        ax = bx + bw
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Inches(0.02),
                                 y0 + Inches(0.6), gap - Inches(0.04),
                                 Inches(0.4))
        arr.fill.solid()
        arr.fill.fore_color.rgb = AMBER
        arr.line.fill.background()

add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
         "Each stage is a standalone binary with a strict, documented contract.",
         size=14, color=DARK)

# Side tools
add_rect(s, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.6), LIGHT)
add_text(s, Inches(0.8), Inches(5.2), Inches(12), Inches(0.4),
         "SIDE TOOLS", size=11, bold=True, color=TEAL)
sides = [
    ("duzip", "CSV <-> zstd"),
    ("duhuman", "machine -> human CSV"),
    ("dumachine", "human -> machine CSV"),
]
for i, (n, d) in enumerate(sides):
    sx = Inches(0.8) + Inches(4.0) * i
    add_text(s, sx, Inches(5.65), Inches(4.0), Inches(0.4),
             n, size=18, bold=True, color=NAVY)
    add_text(s, sx, Inches(6.1), Inches(4.0), Inches(0.4),
             d, size=13, color=GREY)

footer(s)

# ======================================================================
# 3. duscan
# ======================================================================
s = new_slide()
header(s, "duscan", "stage 1 . scanner")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
         "Multi-threaded filesystem walker. Streams POSIX-like metadata for "
         "every file and directory.", size=15, color=DARK)

add_bullets(s, Inches(0.5), Inches(2.6), Inches(6.2), Inches(3.8), [
    ("Workers", " 2 x CPU (cap 48)"),
    ("Batch", " 2048 files per task"),
    ("Flush", " 4 MB threshold"),
    ("Buffer", " 32 MB BufWriter per worker"),
    ("Output", " sharded shards merged into a single file"),
    ("--bin", " zstd binary instead of CSV"),
    ("--no-atime", " zero ATIME for reproducible output"),
], size=15)

add_text(s, Inches(7.0), Inches(2.5), Inches(5.8), Inches(0.4),
         "CSV SCHEMA (9 fields)", size=11, bold=True, color=TEAL)
add_code(s, Inches(7.0), Inches(2.9), Inches(5.8), Inches(1.0),
         "INODE,ATIME,MTIME,UID,GID,\nMODE,SIZE,DISK,PATH", size=14)

add_text(s, Inches(7.0), Inches(4.2), Inches(5.8), Inches(0.4),
         "INVOCATION", size=11, bold=True, color=TEAL)
add_code(s, Inches(7.0), Inches(4.6), Inches(5.8), Inches(1.8),
         "$ duscan /data -o scan.csv\n"
         "$ duscan /data -o scan.zst --bin\n"
         "$ duscan /data -w 32 -f 1.2b\n"
         "$ duscan /a /b -s .cache -vv", size=12)

footer(s)

# ======================================================================
# 4. dusum
# ======================================================================
s = new_slide()
header(s, "dusum", "stage 2 . rollup")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.8),
         "Aggregates raw rows by ancestor folder, owning user, and age bucket. "
         "Every ancestor of every file gets a row, so parent rollups already "
         "include all descendants - no recursive SUM at query time.",
         size=14, color=DARK)

add_text(s, Inches(0.5), Inches(3.3), Inches(6), Inches(0.4),
         "AGE BUCKETS (default: --age 60,600)", size=11, bold=True, color=TEAL)

# Age buckets as colored cards
buckets = [
    ("0", "< 60 d", "recent", RGBColor(0x2E, 0xA0, 0x4A)),
    ("1", "60 - 600 d", "medium", AMBER),
    ("2", ">= 600 d", "old", RGBColor(0xB5, 0x3A, 0x3A)),
]
for i, (num, cond, lbl, col) in enumerate(buckets):
    bx = Inches(0.5) + Inches(2.1) * i
    by = Inches(3.8)
    add_rect(s, bx, by, Inches(2.0), Inches(1.6), WHITE, line=col)
    add_rect(s, bx, by, Inches(2.0), Inches(0.5), col)
    add_text(s, bx, by + Inches(0.05), Inches(2.0), Inches(0.4),
             f"bucket {num}", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, bx, by + Inches(0.65), Inches(2.0), Inches(0.4),
             cond, size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, bx, by + Inches(1.1), Inches(2.0), Inches(0.4),
             lbl, size=12, color=GREY, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.0), Inches(3.3), Inches(6), Inches(0.4),
         "OUTPUT SCHEMA", size=11, bold=True, color=TEAL)
add_code(s, Inches(7.0), Inches(3.8), Inches(5.8), Inches(1.6),
         "path,user,age,\nfiles,size,disk,linked,\naccessed,modified", size=13)

add_text(s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.4),
         "INVOCATION", size=11, bold=True, color=TEAL)
add_code(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.75),
         "$ dusum scan.csv -o scan.sum.csv --age 60,600", size=13)

footer(s)

# ======================================================================
# 5. dudb
# ======================================================================
s = new_slide()
header(s, "dudb", "stage 3 . SQLite ingester")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
         "Offline, one-shot CSV -> SQLite. Never run by the API server.",
         size=15, color=DARK)

add_text(s, Inches(0.5), Inches(2.6), Inches(6), Inches(0.4),
         "WHY (not in-memory anymore)", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(3.0), [
    ("3 GB CSV", " used to expand to 24 GB heap"),
    ("Startup", " now fast and bounded"),
    ("Queries", " remain interactive at 1B files"),
    ("Frontend", " unchanged - API contract preserved"),
], size=14)

add_text(s, Inches(7.0), Inches(2.6), Inches(6), Inches(0.4),
         "INGEST PRAGMAS (safe - DB is rebuildable)", size=11, bold=True,
         color=TEAL)
add_code(s, Inches(7.0), Inches(3.0), Inches(5.8), Inches(1.8),
         "journal_mode = WAL\nsynchronous  = OFF\n"
         "temp_store   = MEMORY\ncache_size   = -262144\n"
         "foreign_keys = OFF", size=12)

add_text(s, Inches(7.0), Inches(5.0), Inches(6), Inches(0.4),
         "INVOCATION", size=11, bold=True, color=TEAL)
add_code(s, Inches(7.0), Inches(5.4), Inches(5.8), Inches(1.3),
         "$ dudb data.sum.csv -o data.db\n"
         "$ dudb data.sum.csv --rebuild", size=13)

footer(s)

# ======================================================================
# 6. SQLite schema (v2)
# ======================================================================
s = new_slide()
header(s, "SQLite schema v2", "what dudb writes")

add_code(s, Inches(0.5), Inches(2.0), Inches(7.3), Inches(4.7),
         "CREATE TABLE users (\n"
         "  id   INTEGER PRIMARY KEY,\n"
         "  name TEXT NOT NULL UNIQUE\n);\n\n"
         "CREATE TABLE paths (\n"
         "  id        INTEGER PRIMARY KEY,\n"
         "  parent_id INTEGER,\n"
         "  full_path TEXT NOT NULL UNIQUE\n);\n"
         "CREATE INDEX idx_paths_parent\n  ON paths(parent_id);\n\n"
         "CREATE TABLE stats (\n"
         "  path_id, user_id, age,\n"
         "  file_count, file_size, disk_bytes,\n"
         "  linked_size, atime, mtime,\n"
         "  PRIMARY KEY (path_id,user_id,age)\n) WITHOUT ROWID;", size=11)

add_text(s, Inches(8.1), Inches(2.0), Inches(5), Inches(0.4),
         "DESIGN NOTES", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(8.1), Inches(2.45), Inches(4.9), Inches(4.8), [
    ("Byte-for-byte paths", " stored exactly as dusum wrote them"),
    ("Synthetic root", " full_path=\"\" above every platform root"),
    ("Windows-native", " C:\\Users, \\\\srv on Windows DBs"),
    ("WITHOUT ROWID", " PK is the natural clustering"),
    ("schema_version", " validated at duapi startup"),
    ("Size scales with", " folders x users x 3 ages, not files"),
], size=13)

footer(s)

# ======================================================================
# 7. duapi server
# ======================================================================
s = new_slide()
header(s, "duapi", "stage 4 . REST API + SPA host")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
         "Axum server. Read-only r2d2 pool over SQLite. Serves the SPA from "
         "a static directory.", size=14, color=DARK)

add_text(s, Inches(0.5), Inches(2.7), Inches(6), Inches(0.4),
         "STARTUP CHECKS", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(3.1), Inches(6), Inches(3.6), [
    ("JWT_SECRET", " required or exits"),
    ("Pool size", " max(num_cpus, 4)"),
    ("Per conn", " query_only + 30 GB mmap + 64 MB cache"),
    ("Schema", " rejects if metadata.schema_version != 2"),
    ("User list", " cached into OnceLock at boot"),
], size=14)

add_text(s, Inches(7.0), Inches(2.7), Inches(6), Inches(0.4),
         "MIDDLEWARE STACK", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(7.0), Inches(3.1), Inches(6), Inches(3.6), [
    ("CORS", " CORS_ORIGIN if set"),
    ("Timeout", " REQUEST_TIMEOUT_SECS (30)"),
    ("Body limit", " MAX_BODY_BYTES (64 KB)"),
    ("Shutdown", " graceful on SIGTERM / SIGINT"),
    ("DB work", " tokio::task::spawn_blocking"),
], size=14)

footer(s)

# ======================================================================
# 8. REST API
# ======================================================================
s = new_slide()
header(s, "REST API", "stage 4 . endpoints")

endpoints = [
    ("GET", "/api/health", "liveness . unauth"),
    ("POST", "/api/login", "username/password -> 24h JWT"),
    ("GET", "/api/users", "admins: all . others: self"),
    ("GET", "/api/folders", "children by user + age (from SQLite)"),
    ("GET", "/api/files", "regular files in dir (live FS, not DB)"),
]
y = Inches(2.1)
for method, route, desc in endpoints:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), WHITE, line=LIGHT)
    col = TEAL if method == "GET" else AMBER
    add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.7), col)
    add_text(s, Inches(0.5), y + Inches(0.17), Inches(0.9), Inches(0.4),
             method, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.55), y + Inches(0.17), Inches(4.0), Inches(0.4),
             route, size=15, bold=True, color=NAVY,
             font="Consolas")
    add_text(s, Inches(5.6), y + Inches(0.2), Inches(7.0), Inches(0.4),
             desc, size=13, color=GREY)
    y += Inches(0.85)

add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
         "/folders and /files accept: path (OS-native), users (CSV), age (0/1/2). "
         "Non-admins must pass their own username.",
         size=12, color=DARK)

footer(s)

# ======================================================================
# 9. Authentication
# ======================================================================
s = new_slide()
header(s, "Authentication", "JWT + OS credentials")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5),
         "Login hits the host OS. On success, issue a 24h HS256 JWT.",
         size=14, color=DARK)

plats = [
    ("macOS", "dscl . -authonly",
     "native directory services"),
    ("Linux", "su <user> -c true",
     "password on stdin; exit code decides"),
    ("Windows", "fake auth (dev only)",
     "matches %USERNAME% / FAKE_USER"),
]
y = Inches(2.8)
for os_name, cmd, note in plats:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.9), LIGHT)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.9), TEAL)
    add_text(s, Inches(0.8), y + Inches(0.12), Inches(2.0), Inches(0.4),
             os_name, size=17, bold=True, color=NAVY)
    add_text(s, Inches(3.0), y + Inches(0.15), Inches(4.5), Inches(0.4),
             cmd, size=13, bold=True, color=DARK, font="Consolas")
    add_text(s, Inches(3.0), y + Inches(0.5), Inches(9), Inches(0.4),
             note, size=12, color=GREY)
    y += Inches(1.0)

add_rect(s, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.1),
         RGBColor(0xFF, 0xF2, 0xE0), line=AMBER)
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
         "WARNING  .  ADMIN_PASSWORD override", size=12, bold=True,
         color=AMBER)
add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
         "Bypasses auth for CI / dev. Grants admin without ADMIN_GROUP. "
         "Never set in production.",
         size=12, color=DARK)

footer(s)

# ======================================================================
# 10. Path normalization
# ======================================================================
s = new_slide()
header(s, "Path normalization", "one contract for all platforms")

add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.8),
         "The DB stores paths byte-for-byte as dusum wrote them. duapi "
         "normalizes the request into that same OS-native form - no second "
         "canonicalization pass.", size=14, color=DARK)

rules = [
    ("\"\"",                 "synthetic root (lists platform roots)"),
    ("\"/\"",                "Unix root"),
    ("\"C:\" / \"C:\\\\\"",  "-> C:\\"),
    ("\"\\\\\\\\srv\\\\s\"", "UNC preserved"),
    ("\"/a//b/./c\"",        "-> /a/b/c  (dup sep + dot collapsed)"),
    ("\"/a/../b\"",          "rejected (400)"),
    ("\"/a\\0/b\"",          "rejected (400 - NUL byte)"),
]
y = Inches(3.0)
for inp, out in rules:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5),
             WHITE if (int((y - Inches(3.0)) / Inches(0.5)) % 2 == 0) else LIGHT)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(3.5), Inches(0.4),
             inp, size=13, color=NAVY, font="Consolas", bold=True)
    add_text(s, Inches(4.5), y + Inches(0.08), Inches(8.0), Inches(0.4),
             out, size=13, color=DARK)
    y += Inches(0.5)

footer(s)

# ======================================================================
# 11. Frontend
# ======================================================================
s = new_slide()
header(s, "Frontend", "SvelteKit 2 + Svelte 5")

add_text(s, Inches(0.5), Inches(2.0), Inches(6), Inches(0.4),
         "STACK", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(2.4), Inches(6), Inches(3.0), [
    ("SvelteKit", " 2 + adapter-static (no SSR)"),
    ("Svelte", " 5 (runes, $state)"),
    ("Tailwind", " 4"),
    ("Build tool", " Vite 6"),
    ("Cache", " idb 8 (1 min TTL per request)"),
], size=14)

add_text(s, Inches(0.5), Inches(5.0), Inches(6), Inches(0.4),
         "DESKTOP VARIANT", size=11, bold=True, color=TEAL)
add_text(s, Inches(0.5), Inches(5.45), Inches(6), Inches(1.6),
         "desktop/ bundles the same UI in Tauri 2 with a Rust backend "
         "(src-tauri/).", size=13, color=DARK)

add_text(s, Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
         "KEY LIB COMPONENTS (browser/src/lib/)", size=11, bold=True,
         color=TEAL)
add_code(s, Inches(7.0), Inches(2.45), Inches(5.8), Inches(4.5),
         "ActionBar          PageNav\n"
         "AgeFilter          PathStats\n"
         "CopyToast          PickerButton\n"
         "FileBar            PickerWrapper\n"
         "FolderBar          SortDropdown\n"
         "Login              Tooltip\n"
         "TreeMap", size=13)

footer(s)

# ======================================================================
# 12. Config
# ======================================================================
s = new_slide()
header(s, "Configuration", "env vars . CLI flags win")

cfg = [
    ("JWT_SECRET",          "(required)", "HMAC secret"),
    ("ADMIN_GROUP",         "(empty)",    "CSV of admin usernames"),
    ("ADMIN_PASSWORD",      "(unset)",    "DEV ONLY override"),
    ("PORT",                "8080",       "listen port"),
    ("STATIC_DIR",          "./public",   "SPA directory"),
    ("CORS_ORIGIN",         "(none)",     "explicit CORS origin"),
    ("TLS_CERT / TLS_KEY",  "(none)",     "enable HTTPS"),
    ("REQUEST_TIMEOUT_SECS","30",         "per-request timeout"),
    ("MAX_BODY_BYTES",      "65536",      "request body cap"),
    ("MAX_PAGE_SIZE",       "2000",       "cap on /folders and /files"),
    ("FAKE_USER",           "%USERNAME%", "Windows dev-auth user"),
]

# Header row
hy = Inches(2.0)
add_rect(s, Inches(0.5), hy, Inches(12.3), Inches(0.45), NAVY)
add_text(s, Inches(0.7), hy + Inches(0.08), Inches(4.5), Inches(0.4),
         "ENV VAR", size=12, bold=True, color=WHITE)
add_text(s, Inches(5.3), hy + Inches(0.08), Inches(2.5), Inches(0.4),
         "DEFAULT", size=12, bold=True, color=WHITE)
add_text(s, Inches(8.0), hy + Inches(0.08), Inches(5), Inches(0.4),
         "PURPOSE", size=12, bold=True, color=WHITE)

y = Inches(2.45)
for i, (k, d, p) in enumerate(cfg):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.36),
             WHITE if i % 2 == 0 else LIGHT)
    add_text(s, Inches(0.7), y + Inches(0.05), Inches(4.5), Inches(0.3),
             k, size=12, bold=True, color=NAVY, font="Consolas")
    add_text(s, Inches(5.3), y + Inches(0.05), Inches(2.5), Inches(0.3),
             d, size=12, color=GREY, font="Consolas")
    add_text(s, Inches(8.0), y + Inches(0.05), Inches(5), Inches(0.3),
             p, size=12, color=DARK)
    y += Inches(0.36)

footer(s)

# ======================================================================
# 13. Quickstart
# ======================================================================
s = new_slide()
header(s, "Quickstart", "build . scan . serve")

add_code(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.5),
         "# 1) Build\n"
         "$ cd rs && cargo build --release\n"
         "$ cd ../browser && npm install && npm run build\n\n"
         "# 2) Scan, roll up, build the DB\n"
         "$ ./rs/target/release/duscan  /data           -o /tmp/data.csv\n"
         "$ ./rs/target/release/dusum   /tmp/data.csv   -o /tmp/data.sum.csv\n"
         "$ ./rs/target/release/dudb    /tmp/data.sum.csv -o /tmp/data.db\n\n"
         "# 3) Serve API + UI\n"
         "$ JWT_SECRET=<secret> ADMIN_GROUP=root,alice \\\n"
         "    ./rs/target/release/duapi /tmp/data.db \\\n"
         "    --static-dir ./browser/build --port 8000", size=13)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         "Open http://localhost:8000  .  log in with an OS account.",
         size=14, color=TEAL, bold=True)

footer(s)

# ======================================================================
# 14. Scale + deployment
# ======================================================================
s = new_slide()
header(s, "Scale + deployment", "what to know before you ship")

add_text(s, Inches(0.5), Inches(2.0), Inches(6), Inches(0.4),
         "SCALE CHARACTERISTICS", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(0.5), Inches(2.45), Inches(6), Inches(4.5), [
    ("DB size", " folders x users x 3 ages (not files)"),
    ("Hot query", " one PK lookup + one range scan"),
    ("Warm cache", " sub-ms per /folders call"),
    ("RSS", " dominated by mmap / page cache"),
    ("Ingest bottleneck", " dudb path cache (swap for DB-backed on extremes)"),
], size=13)

add_text(s, Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
         "DEPLOYMENT", size=11, bold=True, color=TEAL)
add_bullets(s, Inches(7.0), Inches(2.45), Inches(6), Inches(4.5), [
    ("Docker", " multi-stage; dudb as init job, not on start"),
    ("systemd", " unit runs duapi; env file holds JWT_SECRET"),
    ("Proxy", " nginx/caddy terminates TLS; forward to 127.0.0.1"),
    ("Hardening", " non-root user, TLS on, no ADMIN_PASSWORD"),
    ("CI", " cargo test -r + npm run build on master"),
], size=13)

footer(s)

# ======================================================================
# 15. Closing
# ======================================================================
s = new_slide()
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.3), Inches(0.3), Inches(1.0), AMBER)
add_text(s, Inches(0.7), Inches(2.8), Inches(12), Inches(1.2),
         "Questions?", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.6),
         "Full reference: doc/handbook.md", size=20, color=LIGHT)
add_text(s, Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
         "Repo: github.com/sganis/dutopia", size=16, color=TEAL)
add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
         "Dutopia  .  Rust + SvelteKit  .  v4", size=11, color=TEAL,
         bold=True)

# ======================================================================
# Fix total_slides on all slides that used the placeholder
# ======================================================================
total_slides = slide_num[0]
# Re-stamp page numbers
for idx, slide in enumerate(prs.slides, start=1):
    for shp in slide.shapes:
        if shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if " / " in run.text and run.text.strip().endswith(
                            f" / {total_slides - (total_slides - total_slides)}"):
                        pass
# simpler: replace any "X / 14" patterns written at render time
# (we used total_slides=14 initially - update)
OLD = f" / 14"
NEW = f" / {total_slides}"
for slide in prs.slides:
    for shp in slide.shapes:
        if shp.has_text_frame:
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.endswith(OLD):
                        run.text = run.text[:-len(OLD)] + NEW

out = "doc/dutopia-handbook.pptx"
prs.save(out)
print(f"wrote {out}  ({total_slides} slides)")
